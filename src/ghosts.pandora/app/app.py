#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Copyright 2017 Carnegie Mellon University. All Rights Reserved. See LICENSE.md file for terms.

Usage::
    python server.py [<port>]

Send a GET request::
    curl http://localhost

Send a HEAD request::
    curl -I http://localhost

Send a POST request::
    curl -d "foo=bar&bin=baz" http://localhost

"""

import configparser as cp
import os
import random
import re
import ssl
import subprocess as sp
import urllib.parse
import uuid
from http.server import BaseHTTPRequestHandler, HTTPServer
from io import BytesIO
import traceback
import zipstream
from docx import Document
from faker import Faker
from openpyxl import Workbook
from openpyxl.styles import Font
from PIL import Image
from pptx import Presentation
from pptx.util import Inches as pptx_Inches

VERSION = "0.5.25"


class S(BaseHTTPRequestHandler):
    def log_message(self, format, *args):
        if self.headers:
            print(f'{args} - {self.handler} - {self.headers["User-Agent"]}')
        else:
            print(f"{args}")

    fake = Faker()
    image_array = ["png", "gif", "jpg", "jpeg", "pdf", "ico"]
    strict_image_array = ["png", "gif", "jpg", "jpeg"]
    request_url = ""
    file_requested = ""
    handler = "html"

    config = cp.ConfigParser()
    config.read_file(open(r"./config/app.config"))
    payloads = config["payloads"]

    def setup(self):
        BaseHTTPRequestHandler.setup(self)
        self.request.settimeout(30)

    def do_GET(self):
        self.serve_handler()

    def do_HEAD(self):
        self.serve_handler()

    def do_POST(self):
        self.serve_handler()

    def do_PUT(self):
        self.serve_handler()

    def do_DELETE(self):
        self.serve_handler()

    def do_PATCH(self):
        self.serve_handler()

    def serve_handler(self):
        try:
            self.serve_response()
        except Exception as e:
            self.send_error(500, "Internal Server Error")
            self.log_error("Error processing request: %s", str(e))
            self.log_error(traceback.format_exc())

    def send_non_200(self):
        p = self.config.get("non_200s", "percent_is_302")
        if random.randint(2, 100) > (100 - int(p)):
            self.send_response(302)
            x = re.sub(
                "[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}",
                str(uuid.uuid4()),
                self.request_url,
            )
            self.send_header("Location", x)
            print(self.request_url)
            print(x)
            self.end_headers()
        else:
            self.send_response(404)
            self.end_headers()

    def send_standard_headers(self, content_type, file_name=None):
        self.send_response(200)
        server = {self.config.get("headers", "server")}
        self.server_version = f"Server: {server}".encode("utf8")
        if content_type is not None:
            self.send_header("Content-type", content_type)
        if file_name is not None:
            self.send_header(
                "Content-Disposition", f'attachment; filename="{file_name}"'
            )
        if self.file_requested.endswith(".pdf"):
            self.send_header(
                "Content-Disposition", f'attachment; filename="{self.file_requested}"'
            )
        self.end_headers()

    def return_zip(self, file_name):
        self.handler = "zip"
        self.send_standard_headers("application/zip", file_name)
        buf = BytesIO(self.fake.binary(length=random.choice(range(1000, 80000))))
        buf.seek(0, 0)

        with zipstream.ZipFile(mode="w", compression=zipstream.ZIP_DEFLATED) as z:
            folder_name = self.fake.word()
            for _ in range(1, random.randint(1, 47)):
                z.write_iter(
                    f"{folder_name}/{self.fake.word()}.{self.fake.file_extension()}",
                    buf,
                )
            for chunk in z:
                self.wfile.write(chunk)

    def return_binary(self, file_name):
        self.handler = "binary"
        self.send_standard_headers("application/octet-stream", file_name)
        buf = BytesIO(self.fake.binary(length=random.choice(range(1000, 300000))))
        buf.seek(0, 0)
        self.wfile.write(buf.read())

    def return_json(self):
        self.handler = "json"
        self.send_standard_headers("application/json")
        body = self.fake.json(
            data_columns={"Candidates": ["name", "name", "name"]},
            num_rows=random.randint(1, 1000),
        )
        self.wfile.write(body.encode("utf8"))

    def return_csv(self):
        self.handler = "csv"
        self.send_standard_headers("text/csv")
        body = self.fake.csv(
            header=("Name", "Address", "Password"),
            data_columns=("{{name}}", "{{address}}", "{{password}}"),
            num_rows=random.randint(1, 100),
            include_row_ids=True,
        )
        self.wfile.write(body.encode("utf8"))

    def return_text(self, file_name):
        self.handler = "text"
        self.send_standard_headers("text/text")
        body = self.fake.paragraph(nb_sentences=random.randint(1, 13))
        self.wfile.write(body.encode("utf8"))

    def return_stylesheet(self):
        self.handler = "css"
        self.send_standard_headers("text/css")
        fonts = self.config.get("css", "fonts_array").split(",")
        body = (
            "* {font-family:"
            + random.choice(fonts)
            + "} h1 {font-family:"
            + random.choice(fonts)
            + "} body {width:"
            + str(random.randint(65, 100))
            + "%;}"
        )
        self.wfile.write(body.encode("utf8"))

    def return_script(self):
        self.handler = "js"
        self.send_standard_headers("text/javascript")
        body = f"console.log('{self.fake.word()}, /{self.fake.date()}');"
        self.wfile.write(body.encode("utf8"))

    def return_image(self, request_type: str) -> None:
        """
        Generates and returns an image based on the specified request type.

        Parameters:
            request_type (str): The type of image to generate.
                                Accepted values are 'jpg', 'png', 'gif', and 'pdf'.

        Sends the generated image in the HTTP response with the appropriate content type.
        If the height or width parameters are provided in the query string,
        they are used to set the dimensions of the generated image.
        If invalid parameters are provided, a 400 error is sent in the response.
        """
        self.handler = "img"

        # Generate default random dimensions
        height: int = random.randint(10, 2000)
        width: int = random.randint(10, 2000)

        # Parse query parameters for height and width
        qs = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
        try:
            if "h" in qs and qs["h"]:
                height = int(qs["h"][0])
            if "w" in qs and qs["w"]:
                width = int(qs["w"][0])
        except (ValueError, IndexError):
            # Handle invalid height or width values gracefully
            self.send_error(400, "Invalid height or width value.")
            return

        # Determine the appropriate content type
        content_type: str = (
            f"image/{request_type}"
            if request_type != "pdf"
            else f"application/{request_type}"
        )

        # Send standard headers
        self.send_standard_headers(content_type)

        # Create a new image with the specified size and color
        img: Image.Image = Image.new(
            mode="RGB", size=(width, height), color=self.fake.hex_color()
        )

        # Use BytesIO to handle image in memory
        with BytesIO() as buf:
            img_format: str = "JPEG" if request_type == "jpg" else request_type.upper()
            img.save(buf, img_format)
            buf.seek(0)  # Move to the start of the BytesIO buffer
            self.wfile.write(buf.read())

    def return_onenote(self, file_name):
        self.handler = "onenote"
        self.send_standard_headers("application/onenote", file_name)
        buf = BytesIO(self.fake.binary(length=random.choice(range(1000, 300000))))
        buf.seek(0, 0)
        self.wfile.write(buf.read())

    def return_doc_file(self, file_name):
        self.handler = "doc"
        self.send_standard_headers("application/msword", file_name)

        document = Document()
        document.add_heading(self.fake.sentence(), 0)

        p = document.add_paragraph(f"{self.fake.catch_phrase()} ")
        p.add_run(self.fake.word()).bold = True
        p.add_run(f" {self.fake.catch_phrase()} ")
        p.add_run(self.fake.word()).italic = True

        document.add_heading(self.fake.word(), level=1)
        document.add_paragraph(self.fake.catch_phrase(), style="Intense Quote")

        document.add_paragraph(self.fake.catch_phrase(), style="List Bullet")
        document.add_paragraph(self.fake.catch_phrase(), style="List Number")

        rows = 15
        cols = 3
        table = document.add_table(rows=rows, cols=cols)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = "Name"
        hdr_cells[1].text = "Address"
        hdr_cells[2].text = "Email"
        for i in range(1, rows):
            row_cells = table.add_row().cells
            row_cells[0].text = self.fake.name()
            row_cells[1].text = self.fake.address()
            row_cells[2].text = self.fake.email()

        document.add_page_break()

        buf = BytesIO()
        document.save(buf)

        buf.seek(0, 0)
        self.wfile.write(buf.read())

    def return_ppt_file(self, file_name):
        self.handler = "ppt"
        self.send_standard_headers("application/vnd.ms-powerpoint", file_name)
        ppt_file = Presentation()

        Layout = ppt_file.slide_layouts[0]
        first_slide = ppt_file.slides.add_slide(Layout)

        first_slide.shapes.title.text = self.fake.catch_phrase()
        first_slide.placeholders[1].text = self.fake.catch_phrase()

        Second_Layout = ppt_file.slide_layouts[5]
        second_slide = ppt_file.slides.add_slide(Second_Layout)
        second_slide.shapes.title.text = self.fake.catch_phrase()

        textbox = second_slide.shapes.add_textbox(
            pptx_Inches(3), pptx_Inches(1.5), pptx_Inches(3), pptx_Inches(1)
        )
        text_frame = textbox.text_frame
        paragraph = text_frame.add_paragraph()
        paragraph.text = self.fake.catch_phrase()

        buf = BytesIO()
        ppt_file.save(buf)

        buf.seek(0, 0)
        self.wfile.write(buf.read())

    def return_xls_file(self, file_name):
        self.handler = "xls"
        self.send_standard_headers("application/vnd.ms-excel", file_name)
        wb = Workbook()
        wb["Sheet"].sheet_properties.tabColor = self.fake.hex_color().replace("#", "")
        ws0 = wb["Sheet"]
        ws0["A1"] = self.fake.catch_phrase()
        ws0["A1"].font = Font(
            name="Calibri",
            size=40,
            bold=True,
            italic=False,
            strike=False,
            underline="none",
            color=self.fake.hex_color().replace("#", ""),
        )
        for i in range(1, 5):
            wb.create_sheet(f"Sheet_{i}")
            wb[f"Sheet_{i}"].sheet_properties.tabColor = self.fake.hex_color().replace(
                "#", ""
            )
            ws1 = wb[f"Sheet_{i}"]
            for i in range(1, 100):
                for c in range(1, 100):
                    ws1.cell(i, c).value = random.randint(0, 32000)

        buf = BytesIO()
        wb.save(buf)

        buf.seek(0, 0)
        self.wfile.write(buf.read())

    def return_mp4(self, file_name):
        self.handler = "mp4"
        test_movie = "static/test_movie.mp4"
        if not os.path.isfile(test_movie):
            sp.call(
                f"ffmpeg -y -f lavfi -i testsrc=size=1920x1080:rate=1 -vf hue=s=0 -vcodec libx264 -preset superfast -tune zerolatency -pix_fmt yuv420p -t 1000 -movflags +faststart {test_movie}",
                shell=True,
            )

            # Add zero padding to the end of the file https://stackoverflow.com/questions/12768026/write-zeros-to-file-blocks
            stat = os.stat(test_movie)
            with open(test_movie, "r+") as of:
                of.seek(0, os.SEEK_END)
                of.write("\0" * (50 * 1024 * 1024 - stat.st_size))
                of.flush()

        f = open(test_movie, "rb")
        self.send_standard_headers("video/mp4")
        content = f.read()
        self.wfile.write(content)
        f.close()
        return

    def serve_response(self):
        o = urllib.parse.urlparse(self.path)
        self.request_url = o.path
        self.file_requested = os.path.basename(self.request_url)

        # determinations are made in order of priority
        # 1. specific payloads
        # 2. specific file types (.docx, .pptx, .xlsx, .mp4, etc.)
        # 3. specific directory types (css, slides, video, etc.)
        # 4. html

        # 1. payloads: check first in case we are serving a pre-defined url
        for i in self.payloads:
            payload = self.config.get("payloads", i)
            payload_url = payload.split(",")[0]
            payload_file = payload.split(",")[1]
            payload_header = payload.split(",")[2]
            if o.path.startswith(payload_url):
                print(
                    f"{o.path} starts with {payload_url} - dropping {payload_file} with headers of {payload_header}"
                )
                # need to be able to serve bad documents based on url or sequence
                # files = [f for f in os.listdir("./payloads") if os.path.isfile(join("./payloads", f))]
                f = open(f"./payloads/{payload_file}", "rb")
                self.send_standard_headers(payload_header)
                self.wfile.write(f.read())
                f.close()
                return

        # 2. handle specific file extensions
        # https://en.wikipedia.org/wiki/List_of_file_formats
        request_type = None
        if "." in self.request_url:
            request_type = self.request_url.split(".")[-1]

            if request_type in ["doc", "docx", "dotx", "dot", "docm", "dotm", "odt"]:
                self.return_doc_file(self.file_requested)
                return

            elif request_type in ["mp4"]:
                self.return_mp4(self.file_requested)
                return

            elif request_type in [
                "xls",
                "xlsx",
                "xlsm",
                "xlsb",
                "xltm",
                "xla",
                "xlam",
                "xla",
                "ods",
            ]:
                self.return_xls_file(self.file_requested)
                return

            elif request_type in [
                "ppt",
                "pptx",
                "potx",
                "pot",
                "ppsx",
                "pps",
                "pptm",
                "potm",
                "ppsm",
                "odp",
            ]:
                self.return_ppt_file(self.file_requested)
                return

            elif request_type in ["one"]:
                self.return_onenote(self.file_requested)
                return

            elif request_type in ["png", "gif", "jpg", "jpeg", "pdf", "ico"]:
                self.return_image(request_type)
                return

            elif request_type in ["csv"]:
                self.return_csv()
                return

            elif request_type in ["json"]:
                self.return_json()
                return

            elif request_type in ["css"]:
                self.return_stylesheet()
                return

            elif request_type in ["js"]:
                self.return_script()
                return

            elif request_type in ["zip"]:
                self.return_zip(self.file_requested)
                return

            elif request_type in [
                "msi",
                "tar",
                "gz",
                "iso",
                "rar",
                "exe",
                "bin",
                "chm",
            ]:
                self.return_binary(self.file_requested)
                return

            elif request_type in ["txt"]:
                self.return_text(self.file_requested)
                return

        # 3. handle specific directories
        if o.path == ("/about"):
            self.handler = "about"
            self.send_standard_headers("text/html")
            read_me = ""
            with open("../readme.md", "r") as file:
                read_me = file.read()
            body = f"""<html><body><pre>
                GHOSTS PANDORA, version {VERSION}
                Copyright 2022 Carnegie Mellon University. All Rights Reserved. See LICENSE.md file for terms.
                Part of the GHOSTS NPC Orchestration Platform - please email ddupdyke[at]sei.cmu.edu with bugs/requests/other.
                </pre><hr/>
                <h2>README</h2>
                <pre>
                {read_me}
            </pre></body></html>"""
            self.wfile.write(body.encode("utf8"))
            return

        elif o.path.startswith("/video"):
            f = open("./static/mp4.html", "rb")
            self.send_standard_headers("text/html")
            content = f.read()
            self.wfile.write(content)
            f.close()
            return

        elif o.path.startswith("/stream"):
            f = open("./static/player.html", "rb")
            self.send_standard_headers("text/html")
            content = f.read()
            self.wfile.write(content)
            f.close()
            return

        elif o.path.startswith("/css") or o.path.startswith("/styles"):
            self.return_stylesheet()
            return

        elif o.path.startswith("/json"):
            self.return_json()
            return

        elif o.path.startswith("/js") or o.path.startswith("/scripts"):
            self.return_script()
            return

        elif o.path.startswith("/api"):
            self.return_json()
            return

        elif o.path.startswith("/csv"):
            self.return_csv()
            return

        elif o.path.startswith("/pdf"):
            self.return_image("pdf")
            return

        elif o.path.startswith("/onenote") or o.path.startswith("/notebook"):
            self.file_requested = f"{self.fake.word()}.one"
            self.return_onenote(self.file_requested)
            return

        elif o.path.startswith("/docs"):
            self.file_requested = f"{self.fake.word()}.docx"
            self.return_doc_file(self.file_requested)
            return

        elif o.path.startswith("/slides"):
            self.file_requested = f"{self.fake.word()}.pptx"
            self.return_ppt_file(self.file_requested)
            return

        elif o.path.startswith("/sheets"):
            self.file_requested = f"{self.fake.word()}.xlsx"
            self.return_xls_file(self.file_requested)
            return

        elif (
            o.path.startswith("/i/")
            or o.path.startswith("/img")
            or o.path.startswith("/images")
        ):
            request_type = random.choice(self.strict_image_array)
            self.return_image(request_type)
            return

        if random.randint(2, 100) > 95:
            self.send_non_200()
        else:
            self.send_standard_headers("text/html")
            body = ""
            header = f'<script type="text/javascript" src="/scripts/{self.fake.uuid4()}.js"></script><link rel="stylesheet" href="/css/{self.fake.uuid4()}/{self.fake.word()}.css" type="text/css" />'
            for _ in range(random.randint(1, 20)):
                if random.randint(2, 100) > 55:
                    body = body + f"<h3>{self.fake.sentence().replace('.','')}</h3>"
                    body = (
                        body
                        + f"<p>{self.fake.paragraph(nb_sentences=random.randint(1, 100))}</p>"
                    )
                    if random.randint(1, 100) > 85:
                        body = (
                            body
                            + f"<img src='images/{self.fake.word()}.png?h={random.randint(80, 200)}&w={random.randint(200, 400)}'/>"
                        )
            self.wfile.write(
                f"<html><head>{header}<title>{self.fake.catch_phrase()}</title></head><body><h1>{self.fake.catch_phrase()}</h1>{body}</body></html>".encode(
                    "utf8"
                )
            )


def run(server_class=HTTPServer, handler_class=S, port=80):
    server_address = ("0.0.0.0", port)
    httpd = server_class(server_address, handler_class)
    if port == 443:
        if not os.path.exists("./server.pem"):
            print("You told me to run on 443, but no server.pem file exists! Exiting.")
            exit(1)
        httpd.socket = ssl.wrap_socket(
            httpd.socket, certfile="./server.pem", server_side=True
        )

    print("""
░▄▀▒░█▄█░▄▀▄░▄▀▀░▀█▀░▄▀▀░░▒█▀▄▒▄▀▄░█▄░█░█▀▄░▄▀▄▒█▀▄▒▄▀▄
░▀▄█▒█▒█░▀▄▀▒▄██░▒█▒▒▄██▒░░█▀▒░█▀█░█▒▀█▒█▄▀░▀▄▀░█▀▄░█▀█
""")
    print(
        f"Starting GHOSTS PANDORA server {VERSION} on port {port}...\nReady for requests!"
    )
    httpd.serve_forever()


if __name__ == "__main__":
    from sys import argv

    if len(argv) == 2:
        run(port=int(argv[1]))
    else:
        print("Error - you did not supply a port to run on!")
        exit(1)
