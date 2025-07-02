from io import BytesIO
from fastapi import APIRouter, Request, Response
from pptx import Presentation
from faker import Faker

from config.config import OLLAMA_ENABLED, PPT_MODEL
from utils.ollama import generate_document_with_ollama
from utils.content_manager import ContentManager
from app_logging import setup_logger

router = APIRouter()
logger = setup_logger(__name__)
fake = Faker()
cm = ContentManager(default="index", extension="pptx")


def split_content_to_bullets(content: str) -> str:
    bullets = content.split(". ")
    return "\n".join(f"• {b[:200].rstrip()}..." if len(b) > 200 else f"• {b.strip()}" for b in bullets)


def return_ppt(request: Request) -> Response:
    cm.resolve(request)

    if cm.is_storing():
        if content := cm.load():
            return Response(
                content=content,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename={cm.file_name}"},
            )

    presentation = Presentation()
    title = fake.sentence()
    content = fake.paragraph(nb_sentences=3)

    if OLLAMA_ENABLED:
        try:
            t = generate_document_with_ollama("Provide just the title for a PowerPoint slide.", PPT_MODEL)
            title = t or title
        except Exception as e:
            logger.warning(f"Ollama title gen failed: {e}")

        try:
            c = generate_document_with_ollama(
                f"Provide just the content for a PowerPoint presentation slide with the title {title}", PPT_MODEL
            )
            content = c or content
        except Exception as e:
            logger.warning(f"Ollama content gen failed: {e}")

    bullets = split_content_to_bullets(content)

    # Title Slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = fake.sentence()

    # Content Slide
    content_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    content_slide.shapes.title.text = title
    content_slide.placeholders[1].text = bullets

    buf = BytesIO()
    presentation.save(buf)
    buf.seek(0)

    if cm.is_storing():
        cm.save(buf.getvalue())

    logger.info(f"PowerPoint file generated: {cm.file_name}")

    return Response(
        content=buf.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={cm.file_name}"},
    )


# Register routes
ROUTES = ["/ppt", "/slides"]
for route in ROUTES:
    router.add_api_route(f"{route}", return_ppt, methods=["GET", "POST"], tags=["Presentations"])
    router.add_api_route(f"{route}/{{file_name:path}}", return_ppt, methods=["GET", "POST"], tags=["Presentations"])
